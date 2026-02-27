"""PDF and document conversion logic (PDF↔Word/Excel/PPT/Image/PDF-A/HTML)."""
import os
import subprocess
import uuid
import zipfile
from typing import Callable, List, Optional, Tuple

import fitz
import img2pdf
import openpyxl
import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def pdf_to_word(input_path: str, output_path: str) -> List[str]:
    """Convert PDF to Word. Returns list of temp paths to clean up (empty)."""
    doc = fitz.open(input_path)
    word_doc = Document()
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        if page_num > 0:
            word_doc.add_page_break()
        if text.strip():
            word_doc.add_paragraph(text)
    doc.close()
    word_doc.save(output_path)
    return []


def pdf_to_ppt(
    input_path: str, output_path: str, get_output_path: Callable[[str], str]
) -> Tuple[str, List[str]]:
    """Convert PDF to PowerPoint. Returns (output_path, temp_paths)."""
    temp_images: List[str] = []
    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = get_output_path("png")
        pix.save(img_path)
        temp_images.append(img_path)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    doc.close()
    prs.save(output_path)
    return output_path, temp_images


def pdf_to_excel(input_path: str, output_path: str) -> List[str]:
    """Convert PDF to Excel. Returns temp paths (empty). Raises ValueError if empty."""
    all_tables: List[dict] = []
    all_text_pages: List[dict] = []
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            if tables:
                for table_idx, table in enumerate(tables):
                    all_tables.append(
                        {"page": page_num + 1, "table_idx": table_idx + 1, "data": table}
                    )
            else:
                text = page.extract_text()
                if text:
                    all_text_pages.append({"page": page_num + 1, "text": text})
    if not all_tables and not all_text_pages:
        raise ValueError("PDF appears to be empty or image-based (try OCR first)")
    wb = openpyxl.Workbook()
    if all_tables:
        wb.remove(wb.active)
        for table_info in all_tables:
            sheet_name = f"Page {table_info['page']} Table {table_info['table_idx']}"
            ws = wb.create_sheet(title=sheet_name[:31])
            for row_idx, row in enumerate(table_info["data"], start=1):
                for col_idx, cell_value in enumerate(row, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    if row_idx == 1:
                        cell.font = openpyxl.styles.Font(bold=True)
    else:
        ws = wb.active
        ws.title = "Extracted Text"
        ws["A1"] = "No tables found — raw text extracted"
        ws["A1"].font = openpyxl.styles.Font(bold=True)
        row_num = 2
        for page_info in all_text_pages:
            ws.cell(row=row_num, column=1, value=f"--- Page {page_info['page']} ---")
            ws.cell(row=row_num, column=1).font = openpyxl.styles.Font(bold=True)
            row_num += 1
            for line in page_info["text"].split("\n"):
                if line.strip():
                    ws.cell(row=row_num, column=1, value=line)
                    row_num += 1
            row_num += 1
    wb.save(output_path)
    return []


def pdf_to_jpg(
    input_path: str, dpi: int, get_output_path: Callable[[str], str]
) -> Tuple[Optional[str], Optional[str], List[str]]:
    """
    Convert PDF to JPG. Returns (single_jpg_path_or_none, zip_path_or_none, temp_image_paths).
    If single page: single path set, zip_path None. If multi: zip_path set, single None.
    """
    doc = fitz.open(input_path)
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    images: List[str] = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        pixmap = page.get_pixmap(matrix=matrix)
        img_path = get_output_path("jpg")
        pixmap.save(img_path)
        images.append(img_path)
    page_count = len(doc)
    doc.close()
    if page_count == 1:
        return images[0], None, []
    zip_path = get_output_path("zip")
    with zipfile.ZipFile(zip_path, "w") as zf:
        for i, img in enumerate(images):
            zf.write(img, f"page_{str(i+1).zfill(2)}.jpg")
    for img in images:
        if os.path.exists(img):
            os.remove(img)
    return None, zip_path, []


def pdf_to_pdfa(input_path: str, output_path: str, gs_exe: str) -> List[str]:
    """Convert to PDF/A-2b using Ghostscript. Returns temp paths (empty)."""
    result = subprocess.run(
        [
            gs_exe,
            "-dPDFA=2",
            "-dBATCH",
            "-dNOPAUSE",
            "-dNOOUTERSAVE",
            "-sColorConversionStrategy=UseDeviceIndependentColor",
            "-sDEVICE=pdfwrite",
            "-dPDFACompatibilityPolicy=1",
            f"-sOutputFile={output_path}",
            input_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
        raise RuntimeError(
            f"PDF/A conversion failed: {(result.stderr or result.stdout or '')[:200]}"
        )
    return []


def image_to_pdf(input_path: str, output_path: str) -> List[str]:
    """Convert single image to PDF. Returns temp paths (empty)."""
    with open(output_path, "wb") as f:
        f.write(img2pdf.convert(input_path))
    return []


def html_to_pdf_file(input_path: str, output_path: str) -> List[str]:
    """Convert HTML file to PDF (WeasyPrint or reportlab fallback)."""
    try:
        from weasyprint import HTML
        HTML(filename=input_path).write_pdf(output_path)
    except ImportError:
        raise RuntimeError("weasyprint not installed on server")
    except OSError as e:
        if "pango" in str(e).lower() or "library" in str(e).lower():
            _html_fallback_file(input_path, output_path)
        else:
            raise
    return []


def _html_fallback_file(input_path: str, output_path: str) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from html.parser import HTMLParser

    class SimpleHTMLParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []

        def handle_data(self, data):
            if data.strip():
                self.text.append(data.strip())

    with open(input_path, "r", encoding="utf-8") as f:
        html_content = f.read()
    parser = SimpleHTMLParser()
    parser.feed(html_content)
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for text in parser.text:
        story.append(Paragraph(text, styles["Normal"]))
        story.append(Spacer(1, 0.2 * inch))
    doc.build(story)


def html_to_pdf_url(url: str, output_path: str) -> List[str]:
    """Convert URL to PDF (WeasyPrint or reportlab fallback)."""
    try:
        from weasyprint import HTML
        HTML(url=url).write_pdf(output_path)
    except ImportError:
        raise RuntimeError("weasyprint not installed on server")
    except OSError as e:
        if "pango" in str(e).lower() or "library" in str(e).lower():
            _html_fallback_url(url, output_path)
        else:
            raise
    return []


def _html_fallback_url(url: str, output_path: str) -> None:
    import requests
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from html.parser import HTMLParser

    class SimpleHTMLParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []

        def handle_data(self, data):
            if data.strip():
                self.text.append(data.strip())

    response = requests.get(url, timeout=10)
    response.raise_for_status()
    parser = SimpleHTMLParser()
    parser.feed(response.text)
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for text in parser.text:
        story.append(Paragraph(text, styles["Normal"]))
        story.append(Spacer(1, 0.2 * inch))
    doc.build(story)


def office_docx_to_pdf(input_path: str, output_path: str) -> List[str]:
    """Word to PDF using python-docx + reportlab."""
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch

    doc = Document(input_path)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for para in doc.paragraphs:
        if para.text.strip():
            story.append(Paragraph(para.text, styles["Normal"]))
            story.append(Spacer(1, 0.2 * inch))
    pdf_doc.build(story)
    return []


def office_xlsx_to_pdf(input_path: str, output_path: str) -> List[str]:
    """Excel to PDF using openpyxl + reportlab."""
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.pdfgen import canvas as pdf_canvas

    wb = openpyxl.load_workbook(input_path)
    c = pdf_canvas.Canvas(output_path, pagesize=landscape(letter))
    y_position = 550
    for sheet in wb.worksheets:
        c.drawString(50, y_position, f"Sheet: {sheet.title}")
        y_position -= 20
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(
                [str(cell) if cell is not None else "" for cell in row]
            )
            if row_text.strip():
                c.drawString(50, y_position, row_text[:100])
                y_position -= 15
                if y_position < 50:
                    c.showPage()
                    y_position = 550
        c.showPage()
        y_position = 550
    c.save()
    return []


def office_pptx_to_pdf(input_path: str, output_path: str) -> List[str]:
    """PowerPoint to PDF using pptx + reportlab."""
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas as pdf_canvas

    prs = Presentation(input_path)
    c = pdf_canvas.Canvas(output_path, pagesize=landscape(letter))
    page_width, page_height = landscape(letter)
    for slide_num, slide in enumerate(prs.slides, 1):
        y_position = page_height - 80
        c.setFont("Helvetica-Bold", 18)
        c.drawString(50, y_position, f"Slide {slide_num}")
        y_position -= 40
        c.setFont("Helvetica", 11)
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if hasattr(shape, "text_frame") and shape.text_frame:
                c.setFont("Helvetica-Bold", 14)
                c.drawString(70, y_position, text[:90])
                y_position -= 25
                c.setFont("Helvetica", 11)
            else:
                for line in text.split("\n"):
                    if line.strip():
                        if len(line) > 90:
                            words = line.split()
                            current_line = ""
                            for word in words:
                                if len(current_line + word) < 90:
                                    current_line += word + " "
                                else:
                                    c.drawString(90, y_position, current_line)
                                    y_position -= 18
                                    current_line = word + " "
                            if current_line:
                                c.drawString(90, y_position, current_line)
                                y_position -= 18
                        else:
                            c.drawString(90, y_position, line)
                            y_position -= 18
                    y_position -= 10
                    if y_position < 100:
                        break
        c.showPage()
    c.save()
    return []
